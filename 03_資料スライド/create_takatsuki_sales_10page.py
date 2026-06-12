# -*- coding: utf-8 -*-
"""Takatsuki BASE 営業資料（16:9・全10ページ）PPTX 生成スクリプト

ストーリー構成: Problem → Vision → Solution → Proof → Call To Action
「地域情報誌 × 行政パンフレット × スタートアップらしい洗練さ」をテーマに、
街並み・モックアップなどはすべてオリジナルのベクター図形で描画する。
"""
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR

# ========= ブランドカラー（既存1枚資料から継承） =========
GREEN   = RGBColor(0x1f, 0x5c, 0x3d)   # メイン
GREEN_D = RGBColor(0x14, 0x3a, 0x27)   # 濃いグリーン（背景）
GREEN_M = RGBColor(0x35, 0x6b, 0x4c)   # 中間グリーン
GREEN_L = RGBColor(0x8a, 0xb0, 0x9a)   # 淡いグリーン（イラスト）
TINT    = RGBColor(0xee, 0xf3, 0xee)   # カード地（淡）
TINT2   = RGBColor(0xe6, 0xee, 0xe7)   # アイコン円地
ORANGE  = RGBColor(0xe8, 0x83, 0x3a)   # アクセント
ORANGE_D= RGBColor(0xcf, 0x6f, 0x2a)
PEACH   = RGBColor(0xff, 0xd9, 0xb0)
CREAM   = RGBColor(0xfa, 0xf7, 0xf0)
WHITE   = RGBColor(0xff, 0xff, 0xff)
INK     = RGBColor(0x2c, 0x3a, 0x30)   # 本文
SUB     = RGBColor(0x6f, 0x81, 0x70)   # 補足
LINE    = RGBColor(0xd8, 0xe0, 0xd4)   # 罫線
SKY     = RGBColor(0xdf, 0xe9, 0xe2)   # ダーク背景上の補足文字

FONT = "Hiragino Sans"

# ========= スライド寸法（16:9） =========
SW = 13.333
SH = 7.5

prs = Presentation()
prs.slide_width  = In(SW)
prs.slide_height = In(SH)
BLANK = prs.slide_layouts[6]


# ========= ヘルパー =========
def add_slide():
    s = prs.slides.add_slide(BLANK)
    return s, s.shapes


def no_shadow(sp):
    sp.shadow.inherit = False


def rect(shapes, kind, x, y, w, h, fill=None, line=None, line_w=0.75, radius=None):
    sp = shapes.add_shape(kind, In(x), In(y), In(w), In(h))
    no_shadow(sp)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def _set_spc(run, pts):
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(int(pts * 100)))


def txt(shapes, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        line_spacing=1.2, space_after=0):
    """runs: list[paragraph]; paragraph = list[(text, size, color, bold[, letterspacing])]"""
    tb = shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_before = Pt(0)
        if space_after:
            p.space_after = Pt(space_after)
        for (t, sz, col, bold, *rest) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = FONT
            if rest and rest[0]:
                _set_spc(r, rest[0])
    return tb


def line_seg(shapes, x1, y1, x2, y2, color=LINE, w=1.0, dash=False):
    cn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, In(x1), In(y1), In(x2), In(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    no_shadow(cn)
    return cn


def icon_circle(shapes, cx, cy, d, glyph, gsize, bg=TINT2, fg=INK):
    """中心(cx,cy)に直径dの円＋絵文字グリフ"""
    rect(shapes, MSO_SHAPE.OVAL, cx - d/2, cy - d/2, d, d, fill=bg)
    txt(shapes, cx - d/2, cy - d/2, d, d, [[(glyph, gsize, fg, False)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)


def kicker(shapes, x, y, jp, en):
    """セクション小見出し（日本語 + 英語）。装飾下線は使わない。"""
    txt(shapes, x, y, 10.5, 0.4,
        [[(jp, 13, GREEN, True, 0.5), ("   " + en, 9.5, ORANGE, True, 1.5)]],
        anchor=MSO_ANCHOR.MIDDLE)


def page_footer(shapes, n):
    txt(shapes, 0.55, SH - 0.5, 7.0, 0.35,
        [[("Takatsuki BASE", 8.5, SUB, True), ("  ｜ 高槻のお店・人・イベントをつなぐ地域情報ポータル", 8.5, SUB, False)]],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(shapes, SW - 1.7, SH - 0.5, 1.15, 0.35,
        [[("%02d / 10" % n, 8.5, SUB, True)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def bg_fill(shapes, color):
    rect(shapes, MSO_SHAPE.RECTANGLE, 0, 0, SW, SH, fill=color)


def cityscape(shapes, x, y, w, h, base=GREEN_L, accent=None):
    """オリジナルの街並みシルエット（建物＝角丸長方形の並び）。"""
    if accent is None:
        accent = base
    # 建物データ: (相対x, 高さ比, 幅比)
    blocks = [
        (0.00, 0.55, 0.13), (0.12, 0.80, 0.10), (0.21, 0.45, 0.12),
        (0.32, 0.95, 0.11), (0.42, 0.62, 0.13), (0.54, 1.00, 0.09),
        (0.62, 0.50, 0.12), (0.73, 0.78, 0.11), (0.83, 0.58, 0.10),
        (0.92, 0.88, 0.10),
    ]
    for bx, bh, bw in blocks:
        bw_a = bw * w
        bh_a = bh * h
        bxp = x + bx * w
        byp = y + (h - bh_a)
        rect(shapes, MSO_SHAPE.ROUNDED_RECTANGLE, bxp, byp, bw_a, bh_a,
             fill=base, radius=0.06)
        # 窓（小さな点を2列）
        rows = max(1, int(bh_a / 0.42))
        for r in range(rows):
            wy = byp + 0.18 + r * 0.40
            if wy > y + h - 0.2:
                break
            for cxi in (0.28, 0.62):
                wx = bxp + cxi * bw_a
                rect(shapes, MSO_SHAPE.RECTANGLE, wx, wy, 0.09, 0.13, fill=accent)


# =====================================================================
# 1ページ目：表紙
# =====================================================================
s, sh = add_slide()
bg_fill(sh, GREEN_D)
# 右側の装飾円（重なり）
rect(sh, MSO_SHAPE.OVAL, 9.6, -1.7, 6.5, 6.5, fill=RGBColor(0x1b, 0x47, 0x31))
rect(sh, MSO_SHAPE.OVAL, 11.4, 2.2, 4.2, 4.2, fill=GREEN_M)
# 下部の街並み
cityscape(sh, 0, SH - 1.85, SW, 1.55, base=RGBColor(0x10, 0x2e, 0x1f),
          accent=RGBColor(0x2a, 0x55, 0x3b))
# ロゴ
rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, 0.85, 0.75, 0.95, 0.95, fill=ORANGE, radius=0.3)
txt(sh, 0.85, 0.76, 0.95, 0.95, [[("TB", 20, WHITE, True)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(sh, 1.98, 0.72, 7.5, 1.0, [
    [("Takatsuki BASE", 21, WHITE, True)],
    [("TAKATSUKI LOCAL PORTAL", 9, RGBColor(0xb8, 0xcc, 0xc0), False, 2.0)],
], line_spacing=1.05)
# チップ
rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 2.55, 5.0, 0.62, fill=GREEN_M, radius=0.5)
txt(sh, 0.9, 2.55, 5.0, 0.62, [[("地域共創プロジェクト", 11, PEACH, True, 0.5)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# キャッチコピー
txt(sh, 0.85, 3.35, 9.6, 2.3, [
    [("高槻の", 40, WHITE, True), ("「知りたい」", 40, ORANGE, True)],
    [("と", 40, WHITE, True), ("「伝えたい」", 40, PEACH, True), ("を、", 40, WHITE, True)],
    [("ひとつに。", 40, WHITE, True)],
], line_spacing=1.15)
# サブ
txt(sh, 0.9, SH - 1.25, 11.5, 0.7,
    [[("高槻のお店・企業・イベントをつなぐ、地域情報プラットフォーム", 13.5, SKY, False, 0.3)]])

# =====================================================================
# 2ページ目：高槻の地域情報はもっとつながれる（社会背景）
# =====================================================================
s, sh = add_slide()
bg_fill(sh, WHITE)
kicker(sh, 0.55, 0.55, "高槻の地域情報は、もっとつながれる", "BACKGROUND")
txt(sh, 0.55, 1.05, 12.2, 1.3, [
    [("便利な時代なのに、", 26, INK, True), ("「地元の良いお店」", 26, GREEN, True),
     ("ほど", 26, INK, True)],
    [("見つけにくい", 26, ORANGE, True), ("──そんな状況が生まれています。", 26, INK, True)],
], line_spacing=1.2)

bg_items = [
    ("📱", "SNSだけでは埋もれる", "投稿は流れて消え、過去の良い情報まで届きにくい"),
    ("🗺️", "知る機会が少ない", "近所にどんなお店があるか、意外と知られていない"),
    ("🧳", "新住民が探しづらい", "高槻へ来たばかりの人が、地元情報にたどり着けない"),
    ("📰", "情報が分散している", "イベント・求人・店舗情報がバラバラの場所にある"),
]
cw = 2.92; gap = 0.30; cx0 = 0.55; cy = 3.05; chh = 3.4
for i, (ic, ti, de) in enumerate(bg_items):
    x = cx0 + i * (cw + gap)
    rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, cw, chh, fill=TINT, radius=0.06)
    icon_circle(sh, x + cw/2, cy + 0.95, 1.15, ic, 26, bg=WHITE, fg=INK)
    txt(sh, x + 0.25, cy + 1.7, cw - 0.5, 0.7, [[(ti, 13.5, GREEN, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.1)
    txt(sh, x + 0.3, cy + 2.35, cw - 0.6, 0.95, [[(de, 10.5, SUB, False)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.3)
page_footer(sh, 2)

# =====================================================================
# 3ページ目：地域事業者が抱える課題（PROBLEM）
# =====================================================================
s, sh = add_slide()
bg_fill(sh, WHITE)
kicker(sh, 0.55, 0.55, "地域のお店が抱える、5つの「困った」", "PROBLEM")
txt(sh, 0.55, 1.05, 12.2, 0.8,
    [[("いいお店なのに知られていない──その背景には、共通の悩みがあります。", 15, INK, False)]])

probs = [
    ("🙍", "新規のお客様が増えない", "リピーターは来てくれるが、新しい層に広がらない"),
    ("📣", "SNSだけでは限界がある", "フォロワー以外には、なかなか情報が届かない"),
    ("💻", "HPの制作・更新が負担", "作る時間も、こまめに更新する余裕もない"),
    ("🎪", "告知する場所が少ない", "イベントやセールを、どこで広めればいいか分からない"),
    ("🧑‍🍳", "求人媒体は費用が高い", "募集を出すたびにコストがかさんでしまう"),
]
# 上段3 + 下段2
cw3 = 3.95; g3 = 0.30; x0 = 0.55; rowy = [2.0, 4.5]; chh = 2.25
positions = [(0, 0), (1, 0), (2, 0), (0, 1), (1, 1)]
for idx, (ic, ti, de) in enumerate(probs):
    col, row = positions[idx]
    x = x0 + col * (cw3 + g3)
    y = rowy[row]
    rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw3, chh, fill=WHITE, line=LINE, line_w=1.0, radius=0.07)
    rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 0.12, chh, fill=ORANGE, radius=0.5)
    icon_circle(sh, x + 0.95, y + 0.72, 1.05, ic, 22, bg=TINT, fg=INK)
    txt(sh, x + 1.65, y + 0.32, cw3 - 1.85, 0.8, [[(ti, 13.5, INK, True)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    txt(sh, x + 0.45, y + 1.38, cw3 - 0.8, 0.8, [[(de, 10.5, SUB, False)]],
        line_spacing=1.25)
# 下段右の余白にメッセージ
rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, x0 + 2 * (cw3 + g3), rowy[1], cw3, chh, fill=GREEN, radius=0.07)
txt(sh, x0 + 2 * (cw3 + g3) + 0.45, rowy[1] + 0.35, cw3 - 0.9, chh - 0.7, [
    [("これ、", 14, WHITE, True), ("ひとつでも", 14, PEACH, True)],
    [("当てはまりませんか？", 14, WHITE, True)],
    [("", 6, WHITE, False)],
    [("次のページから、解決の", 10.5, SKY, False)],
    [("かたちをご紹介します。", 10.5, SKY, False)],
], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
page_footer(sh, 3)

# =====================================================================
# 4ページ目：目指したい未来（VISION）
# =====================================================================
s, sh = add_slide()
bg_fill(sh, GREEN)
# 装飾円
rect(sh, MSO_SHAPE.OVAL, -1.6, 3.6, 5.5, 5.5, fill=GREEN_M)
rect(sh, MSO_SHAPE.OVAL, 11.2, -1.8, 4.6, 4.6, fill=GREEN_M)
txt(sh, 0.55, 0.6, 7.0, 0.4, [[("もし、こんな高槻になったら", 13, PEACH, True, 0.5), ("   VISION", 9.5, ORANGE, True, 1.5)]],
    anchor=MSO_ANCHOR.MIDDLE)
txt(sh, 0.55, 1.1, 12.2, 1.7, [
    [("お店・人・イベントが", 30, WHITE, True)],
    [("自然につながる", 30, PEACH, True), ("、高槻へ。", 30, WHITE, True)],
], line_spacing=1.15)

vis = [
    ("🔎", "お店が見つかる", "地域のお店が、必要な人にちゃんと届く"),
    ("🧭", "情報がすぐ探せる", "高槻で暮らす人が、知りたい情報にすぐ会える"),
    ("🤝", "人とお店がつながる", "お店・人・イベントが、自然に結びつく"),
    ("🌱", "街が盛り上がる", "地域全体に、にぎわいの循環が生まれる"),
]
cw = 2.92; gap = 0.30; x0 = 0.55; cy = 3.5; chh = 3.0
for i, (ic, ti, de) in enumerate(vis):
    x = x0 + i * (cw + gap)
    rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, cw, chh, fill=WHITE, radius=0.07)
    icon_circle(sh, x + cw/2, cy + 0.85, 1.1, ic, 25, bg=TINT, fg=INK)
    txt(sh, x + 0.2, cy + 1.55, cw - 0.4, 0.6, [[(ti, 13.5, GREEN, True)]],
        align=PP_ALIGN.CENTER, line_spacing=1.05)
    txt(sh, x + 0.3, cy + 2.1, cw - 0.6, 0.8, [[(de, 10.5, SUB, False)]],
        align=PP_ALIGN.CENTER, line_spacing=1.3)

# =====================================================================
# 5ページ目：Takatsuki BASEとは（SOLUTION）
# =====================================================================
s, sh = add_slide()
bg_fill(sh, WHITE)
kicker(sh, 0.55, 0.55, "その未来をかたちにする仕組みが、これです", "SOLUTION")
# ロゴ + 概要（左）
rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, 1.2, 0.85, 0.85, fill=ORANGE, radius=0.3)
txt(sh, 0.55, 1.21, 0.85, 0.85, [[("TB", 17, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(sh, 1.6, 1.18, 6.0, 0.95, [
    [("Takatsuki BASE", 22, GREEN, True)],
    [("高槻の地域情報ポータルサイト", 11.5, SUB, False)],
], line_spacing=1.1)
txt(sh, 0.55, 2.45, 6.3, 2.2, [
    [("高槻市内の", 14, INK, False), ("飲食店・お店・企業・イベント・求人", 14, GREEN, True),
     ("の情報を、", 14, INK, False)],
    [("ひとつのサイトにまとめた地域情報ポータルです。", 14, INK, False)],
    [("", 7, INK, False)],
    [("「広告を売る」のではなく、", 14, INK, False)],
    [("みんなで高槻を盛り上げる", 14, ORANGE, True), ("ことを目的に運営します。", 14, INK, False)],
], line_spacing=1.45)
# 提供価値3点
vals = [("集める", "地域の情報を1か所に集約"), ("見つける", "住む人・働く人が探せる"), ("つながる", "お店と人を自然に結ぶ")]
vy = 5.1
for i, (t, d) in enumerate(vals):
    x = 0.55 + i * 2.18
    rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, x, vy, 2.0, 1.55, fill=TINT, radius=0.1)
    txt(sh, x, vy + 0.28, 2.0, 0.55, [[(t, 13, GREEN, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(sh, x + 0.15, vy + 0.85, 1.7, 0.6, [[(d, 9.5, SUB, False)]], align=PP_ALIGN.CENTER, line_spacing=1.25)

# 右側：集約のコンセプト図（中心ノード＋周辺ノード）
cxc, cyc = 10.55, 3.95
# 接続線（先に描いて背面に）
nodes = [("飲食店", 10.55, 1.75), ("お店", 12.55, 3.0),
         ("イベント", 12.55, 4.9), ("求人", 10.55, 6.15),
         ("企業", 8.55, 4.9), ("住民", 8.55, 3.0)]
for _, nx, ny in nodes:
    line_seg(sh, cxc, cyc, nx, ny, color=LINE, w=1.5)
# 中心
rect(sh, MSO_SHAPE.OVAL, cxc - 1.05, cyc - 1.05, 2.1, 2.1, fill=GREEN)
txt(sh, cxc - 1.05, cyc - 1.05, 2.1, 2.1, [
    [("Takatsuki", 12.5, WHITE, True)], [("BASE", 12.5, PEACH, True)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
# 周辺
for nm, nx, ny in nodes:
    rect(sh, MSO_SHAPE.OVAL, nx - 0.7, ny - 0.7, 1.4, 1.4, fill=WHITE, line=GREEN_L, line_w=1.5)
    txt(sh, nx - 0.7, ny - 0.7, 1.4, 1.4, [[(nm, 11, GREEN, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
page_footer(sh, 5)

# =====================================================================
# 6ページ目：掲載できる内容（＋モックアップ）
# =====================================================================
s, sh = add_slide()
bg_fill(sh, WHITE)
kicker(sh, 0.55, 0.55, "あなたのお店の「載せたい」を、まるごと", "WHAT YOU CAN POST")
txt(sh, 0.55, 1.05, 8.0, 0.7, [[("お店の魅力を伝える情報を、自由に掲載できます。", 15, INK, False)]])

posts = [
    ("🏪", "店舗紹介", "お店の基本情報・こだわり"),
    ("🍽️", "メニュー・サービス", "商品や提供メニューの紹介"),
    ("📷", "写真ギャラリー", "店内・商品の魅力を写真で"),
    ("🎉", "イベント・キャンペーン", "セールや催しの告知"),
    ("🧑‍💼", "求人募集", "スタッフ募集の掲載"),
    ("🔗", "SNS・HP連携", "Instagram・HPとつなぐ"),
]
cw = 3.85; gap = 0.3; x0 = 0.55; rys = [2.0, 3.55, 5.1]; chh = 1.35
for i, (ic, ti, de) in enumerate(posts):
    col = i % 2; row = i // 2
    x = x0 + col * (cw + gap)
    y = rys[row]
    rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, chh, fill=TINT, radius=0.1)
    icon_circle(sh, x + 0.78, y + chh/2, 0.92, ic, 19, bg=WHITE, fg=INK)
    txt(sh, x + 1.45, y + 0.22, cw - 1.6, 0.55, [[(ti, 12.5, GREEN, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(sh, x + 1.45, y + 0.72, cw - 1.6, 0.5, [[(de, 10, SUB, False)]], anchor=MSO_ANCHOR.MIDDLE)

# 右側：スマホ風モックアップ（架空）
mx, my, mw, mh = 9.5, 1.95, 3.0, 5.05
rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, mx, my, mw, mh, fill=RGBColor(0x22, 0x2b, 0x26), radius=0.09)
rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, mx + 0.18, my + 0.18, mw - 0.36, mh - 0.36, fill=WHITE, radius=0.07)
ix = mx + 0.34; iw = mw - 0.68
# ヘッダー帯
rect(sh, MSO_SHAPE.RECTANGLE, ix, my + 0.34, iw, 0.66, fill=GREEN)
txt(sh, ix, my + 0.34, iw, 0.66, [[("Takatsuki BASE", 10, WHITE, True)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 写真プレースホルダ
rect(sh, MSO_SHAPE.RECTANGLE, ix, my + 1.05, iw, 1.25, fill=GREEN_L)
txt(sh, ix, my + 1.05, iw, 1.25, [[("PHOTO", 11, WHITE, True, 1.5)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 店名・テキスト行
txt(sh, ix + 0.05, my + 2.42, iw - 0.1, 0.4, [[("◯◯ベーカリー（サンプル）", 9.5, INK, True)]])
for k in range(3):
    rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, ix, my + 2.92 + k * 0.32, iw * (0.95 - k * 0.18), 0.16,
         fill=LINE, radius=0.5)
# ボタン
rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, ix, my + 4.05, iw, 0.5, fill=ORANGE, radius=0.3)
txt(sh, ix, my + 4.05, iw, 0.5, [[("お店の詳細を見る", 10, WHITE, True)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(sh, mx, my + mh + 0.05, mw, 0.35, [[("※架空の掲載イメージ", 8.5, SUB, False)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
page_footer(sh, 6)

# =====================================================================
# 7ページ目：掲載するメリット（機能→効果 / PROOF）
# =====================================================================
s, sh = add_slide()
bg_fill(sh, WHITE)
kicker(sh, 0.55, 0.55, "掲載すると、こんな効果が期待できます", "BENEFIT")
txt(sh, 0.55, 1.05, 12.2, 0.7, [[("「機能」ではなく、お店が得られる", 15, INK, False),
    ("「効果」", 15, ORANGE, True), ("でご紹介します。", 15, INK, False)]])

# テーブル風: 機能 → 効果
ben = [
    ("地域への掲載", "地域住民への認知向上", "高槻に住む人・働く人にお店が届く"),
    ("検索エンジン対応", "Google検索からの流入", "「高槻 ◯◯」で見つけてもらえる"),
    ("SNS連携", "SNSとの相乗効果", "投稿とサイトで情報の届く範囲が広がる"),
    ("イベント掲載", "イベント集客", "セール・催しの告知で来店を後押し"),
    ("求人掲載", "採用導線の追加", "求人媒体に頼らず募集を届けられる"),
    ("店舗ページ", "HP代わりに使える", "制作・更新の負担なく情報発信できる"),
]
# ヘッダー行
hy = 1.95; rh = 0.63; x0 = 0.55
col_fx, col_fw = 0.55, 3.0      # 機能
col_ax, col_aw = 3.75, 3.4      # 効果
col_dx, col_dw = 7.35, 5.45     # 説明
rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, x0, hy, 12.23, rh, fill=GREEN, radius=0.08)
txt(sh, col_fx + 0.25, hy, col_fw, rh, [[("できること（機能）", 12, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
txt(sh, col_ax + 0.2, hy, col_aw, rh, [[("得られる効果", 12, PEACH, True)]], anchor=MSO_ANCHOR.MIDDLE)
txt(sh, col_dx + 0.2, hy, col_dw, rh, [[("どういうこと？", 12, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
ry = hy + rh + 0.1
for i, (fn, ef, de) in enumerate(ben):
    y = ry + i * (rh + 0.07)
    rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, x0, y, 12.23, rh, fill=(TINT if i % 2 == 0 else WHITE),
         line=LINE, line_w=0.75, radius=0.08)
    txt(sh, col_fx + 0.25, y, col_fw, rh, [[(fn, 11.5, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(sh, col_ax + 0.2, y, col_aw, rh, [[("▶ ", 11, ORANGE, True), (ef, 11.5, GREEN, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(sh, col_dx + 0.2, y, col_dw, rh, [[(de, 10.5, SUB, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
page_footer(sh, 7)

# =====================================================================
# 8ページ目：利用イメージ・活用シーン
# =====================================================================
s, sh = add_slide()
bg_fill(sh, WHITE)
kicker(sh, 0.55, 0.55, "たとえば、こんなふうに使えます", "USE CASE")
txt(sh, 0.55, 1.05, 12.2, 0.7, [[("掲載したあとの、うれしい「その先」をイメージしてみてください。", 15, INK, False)]])

cases = [
    ("🥐", "新商品が来店に", "新メニューを掲載 → 見た人が「食べてみたい」と来店"),
    ("🎈", "イベントに人が集まる", "セール告知を掲載 → 地域住民が当日に足を運ぶ"),
    ("📨", "求人に応募が届く", "スタッフ募集を掲載 → 近所の人が応募してくれる"),
    ("🧳", "観光客が立ち寄る", "高槻を訪れた人がサイトで見つけ、お店に立ち寄る"),
]
cw = 5.95; gap = 0.33; x0 = 0.55; rys = [2.05, 4.6]; chh = 2.35
for i, (ic, ti, de) in enumerate(cases):
    col = i % 2; row = i // 2
    x = x0 + col * (cw + gap)
    y = rys[row]
    rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, chh, fill=TINT, radius=0.07)
    # 番号
    rect(sh, MSO_SHAPE.OVAL, x + 0.4, y + 0.4, 0.7, 0.7, fill=GREEN)
    txt(sh, x + 0.4, y + 0.4, 0.7, 0.7, [[("%d" % (i + 1), 15, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    icon_circle(sh, x + cw - 1.0, y + chh/2, 1.25, ic, 30, bg=WHITE, fg=INK)
    txt(sh, x + 1.3, y + 0.45, cw - 2.6, 0.6, [[(ti, 15, GREEN, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(sh, x + 0.45, y + 1.25, cw - 2.3, 0.95, [[(de, 11, INK, False)]], line_spacing=1.35)
page_footer(sh, 8)

# =====================================================================
# 9ページ目：Takatsuki BASEが目指すこと（共創・地域貢献）
# =====================================================================
s, sh = add_slide()
bg_fill(sh, GREEN_D)
rect(sh, MSO_SHAPE.OVAL, 10.6, 3.8, 5.5, 5.5, fill=RGBColor(0x1b, 0x47, 0x31))
txt(sh, 0.55, 0.65, 8.0, 0.4, [[("私たちが目指していること", 13, PEACH, True, 0.5), ("   OUR MISSION", 9.5, ORANGE, True, 1.5)]],
    anchor=MSO_ANCHOR.MIDDLE)
txt(sh, 0.55, 1.15, 12.2, 1.6, [
    [("Takatsuki BASE は、", 27, WHITE, True)],
    [("単なる広告媒体では", 27, WHITE, True), ("ありません。", 27, PEACH, True)],
], line_spacing=1.15)

miss = [
    ("🤝", "つなぐ", "地域のお店・企業・人をつなぐ"),
    ("📣", "広める", "高槻の魅力を市内外へ発信する"),
    ("📈", "活性化", "地域経済の活性化に貢献する"),
    ("🌱", "育てる", "事業者と一緒にサービスを育てる"),
]
cw = 2.92; gap = 0.30; x0 = 0.55; cy = 3.4; chh = 2.65
for i, (ic, ti, de) in enumerate(miss):
    x = x0 + i * (cw + gap)
    rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, cw, chh, fill=GREEN, radius=0.08)
    icon_circle(sh, x + cw/2, cy + 0.78, 1.0, ic, 23, bg=GREEN_M, fg=WHITE)
    txt(sh, x + 0.2, cy + 1.4, cw - 0.4, 0.55, [[(ti, 15, PEACH, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(sh, x + 0.25, cy + 1.95, cw - 0.5, 0.65, [[(de, 10.5, SKY, False)]],
        align=PP_ALIGN.CENTER, line_spacing=1.3)
txt(sh, 0.55, cy + chh + 0.35, 12.2, 0.6,
    [[("「共創」", 15, ORANGE, True), ("──地域事業者のみなさんと一緒に、高槻のための仕組みを育てていきます。", 14, WHITE, False)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# 10ページ目：お問い合わせ・次のアクション（CALL TO ACTION）
# =====================================================================
s, sh = add_slide()
bg_fill(sh, GREEN)
rect(sh, MSO_SHAPE.OVAL, -1.8, -1.8, 5.0, 5.0, fill=GREEN_M)
cityscape(sh, 0, SH - 1.3, SW, 1.1, base=RGBColor(0x18, 0x42, 0x2c), accent=GREEN_M)
txt(sh, 0.7, 0.7, 8.0, 0.4, [[("次のアクション", 13, PEACH, True, 0.5), ("   CALL TO ACTION", 9.5, ORANGE, True, 1.5)]],
    anchor=MSO_ANCHOR.MIDDLE)
txt(sh, 0.7, 1.25, 8.5, 2.0, [
    [("まずは", 30, WHITE, True), ("無料", 30, PEACH, True), ("で、", 30, WHITE, True)],
    [("掲載の相談から", 30, WHITE, True)],
    [("始めませんか？", 30, WHITE, True)],
], line_spacing=1.18)
txt(sh, 0.7, 3.95, 7.7, 1.3, [
    [("掲載内容の整理や写真選びも、こちらでサポートします。", 12.5, SKY, False)],
    [("「話だけ聞いてみたい」も大歓迎。お気軽にご相談ください。", 12.5, SKY, False)],
], line_spacing=1.5)
# 連絡先
contacts = [("運営", "Takatsuki BASE 運営事務局"), ("WEB", "takatsuki-base.jp"),
            ("Mail", "info@takatsuki-base.jp"), ("IG", "@takatsuki_base")]
cyy = 5.45
for i, (k, v) in enumerate(contacts):
    col = i % 2; row = i // 2
    x = 0.7 + col * 4.1
    y = cyy + row * 0.62
    txt(sh, x, y, 1.0, 0.5, [[(k, 11, ORANGE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(sh, x + 1.0, y, 3.0, 0.5, [[(v, 11.5, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
# QRカード
qx, qy, qw, qh = 9.55, 1.45, 3.1, 4.0
rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, qx, qy, qw, qh, fill=WHITE, radius=0.08)
try:
    sh.add_picture("Takatsuki_BASE_QR.png", In(qx + (qw - 2.2)/2), In(qy + 0.4), In(2.2), In(2.2))
except Exception:
    rect(sh, MSO_SHAPE.RECTANGLE, qx + (qw - 2.2)/2, qy + 0.4, 2.2, 2.2, fill=TINT)
txt(sh, qx, qy + 2.75, qw, 0.5, [[("スマホで読み取り", 11, GREEN, True)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(sh, qx, qy + 3.25, qw, 0.5, [[("takatsuki-base.jp", 10, SUB, False)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 締めメッセージ
rect(sh, MSO_SHAPE.ROUNDED_RECTANGLE, 9.55, 5.65, 3.1, 1.05, fill=ORANGE, radius=0.15)
txt(sh, 9.65, 5.7, 2.9, 0.95, [
    [("一緒に高槻を", 12.5, WHITE, True)],
    [("盛り上げる仲間を募集中", 11, WHITE, True)],
], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
txt(sh, 0.7, SH - 0.42, 9.0, 0.3, [[("※ロゴ・連絡先・QRは差し替え用のサンプルです", 8.5, RGBColor(0xb8, 0xcc, 0xc0), False)]],
    anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
prs.save("Takatsuki_BASE_営業資料_10ページ.pptx")
print("saved: Takatsuki_BASE_営業資料_10ページ.pptx  (slides=%d)" % len(prs.slides._sldIdLst))
