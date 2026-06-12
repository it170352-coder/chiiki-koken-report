# -*- coding: utf-8 -*-
"""Takatsuki BASE 営業1枚資料（A4縦）PPTX 生成スクリプト
HTML版と同じ5ゾーン構成を、PowerPoint/Canvaで編集できる形で出力する。
"""
import os
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- カラー ----
GREEN   = RGBColor(0x2c, 0x81, 0x59)
GREEN_D = RGBColor(0x1f, 0x63, 0x47)
ORANGE  = RGBColor(0xdb, 0x71, 0x50)   # テラコッタ（明るい背景用）
ACCENT_LT = RGBColor(0xf2, 0xa2, 0x79) # 緑地用の明るめコーラル
CREAM   = RGBColor(0xfb, 0xf8, 0xf2)
INK     = RGBColor(0x2a, 0x3b, 0x31)
SUB     = RGBColor(0x6f, 0x81, 0x70)
LINE    = RGBColor(0xe6, 0xe1, 0xd3)
WHITE   = RGBColor(0xff, 0xff, 0xff)
PEACH   = RGBColor(0xff, 0xd9, 0xb0)
FONT = "Hiragino Sans"

prs = Presentation()
prs.slide_width  = Cm(21)
prs.slide_height = Cm(29.7)
slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = slide.shapes

# 背景クリーム
bg = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Cm(21), Cm(29.7))
bg.fill.solid(); bg.fill.fore_color.rgb = CREAM; bg.line.fill.background()
bg.shadow.inherit = False


def no_shadow(sp):
    sp.shadow.inherit = False


def rect(shape_type, x, y, w, h, fill=None, line=None, line_w=0.75, radius=None):
    sp = shapes.add_shape(shape_type, Cm(x), Cm(y), Cm(w), Cm(h))
    no_shadow(sp)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    if radius is not None and shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def txt(x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        line_spacing=1.1, space_after=0):
    """runs: list of paragraphs; each paragraph = list of (text,size,color,bold,spacing)"""
    tb = shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (t, sz, col, bold, *rest) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold
            r.font.name = FONT
            if rest and rest[0]:
                _set_letter_spacing(r, rest[0])
    return tb


def _set_letter_spacing(run, pts):
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(int(pts * 100)))


# ============ ① ヘッダー ============
rect(MSO_SHAPE.RECTANGLE, 0, 0, 21, 5.4, fill=GREEN)
# 装飾の円（オレンジ半透明っぽく＝薄め）
circ = shapes.add_shape(MSO_SHAPE.OVAL, Cm(16.5), Cm(3.4), Cm(5.5), Cm(5.5))
no_shadow(circ); circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x40, 0x93, 0x6a)
circ.line.fill.background()
# ロゴ角丸
rect(MSO_SHAPE.ROUNDED_RECTANGLE, 1.2, 0.7, 1.3, 1.3, fill=ORANGE, radius=0.3)
txt(1.2, 0.72, 1.3, 1.3, [[("TB", 18, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# ロゴテキスト
txt(2.75, 0.62, 9, 1.4, [
    [("Takatsuki BASE", 18, WHITE, True)],
    [("TAKATSUKI LOCAL PORTAL", 8, RGBColor(0xcf, 0xdd, 0xd4), False, 1.5)],
], line_spacing=1.0)
# リージョンチップ
rect(MSO_SHAPE.ROUNDED_RECTANGLE, 13.0, 0.85, 6.8, 0.85, fill=RGBColor(0x3c, 0x90, 0x67), radius=0.5)
txt(13.0, 0.86, 6.8, 0.85, [[("高槻のお店・人・イベントをつなぐ", 9.5, WHITE, True)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# キャッチ
txt(1.2, 2.3, 18.6, 2.2, [
    [("高槻の", 27, WHITE, True), ("「知りたい」", 27, ACCENT_LT, True), ("と", 27, WHITE, True)],
    [("「伝えたい」", 27, PEACH, True), ("をつなぐ。", 27, WHITE, True)],
], line_spacing=1.2)
txt(1.2, 4.55, 18.6, 0.7, [[("地域のお店・企業・イベント情報をまとめた、高槻のための情報ポータルサイト",
    10.5, RGBColor(0xdf, 0xe9, 0xe2), False)]])

# ============ ② サービス概要 ============
rect(MSO_SHAPE.RECTANGLE, 0, 5.4, 21, 2.05, fill=WHITE)
rect(MSO_SHAPE.RECTANGLE, 0, 7.43, 21, 0.02, fill=LINE)
rect(MSO_SHAPE.ROUNDED_RECTANGLE, 1.2, 5.85, 4.7, 1.05, fill=GREEN, radius=0.25)
txt(1.2, 5.86, 4.7, 1.05, [[("Takatsuki BASE とは？", 13, WHITE, True)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(6.3, 5.6, 13.5, 1.8, [
    [("高槻市内の", 13, INK, False), ("飲食店・お店・企業・イベント・求人", 13, GREEN, True),
     ("の情報を1か所にまとめた地域情報ポータルサイトです。", 13, INK, False)],
    [("「広告を売る」のではなく、", 13, INK, False), ("みんなで高槻を盛り上げる", 13, GREEN, True),
     ("ことを目的に運営しています。", 13, INK, False)],
], line_spacing=1.35, anchor=MSO_ANCHOR.MIDDLE)

# ============ ③ お悩み（共感・現場の声） ============
txt(1.2, 7.75, 16, 0.6, [[("こんなお悩み、ありませんか？", 12, GREEN, True, 1.0), ("   YOUR VOICE", 9, ORANGE, True, 1.5)]])
rect(MSO_SHAPE.RECTANGLE, 1.2, 8.4, 18.6, 0.02, fill=LINE)

ICON_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "icons")
problems = [
    ("p_store",   "「いいお店なのに、まだ十分に知られていない気がする」"),
    ("p_trend",   "「チラシやSNSを頑張っても、新規のお客様が増えない」"),
    ("p_search",  "「『高槻 ◯◯』で検索しても、うちのお店が出てこない」"),
    ("p_clock",   "「ホームページを作る時間も、更新する余裕もない」"),
    ("p_mega",    "「イベントやセールの告知、どう広めればいいか分からない」"),
    ("p_usearch", "「求人を出しても、募集していること自体が伝わらない」"),
]
pcols = 3
pgap = 0.4
pw = (18.6 - pgap * (pcols - 1)) / pcols
ph_card = 2.35
px0 = 1.2
py0 = 8.7
for i, (ico, text) in enumerate(problems):
    col = i % pcols
    row = i // pcols
    x = px0 + col * (pw + pgap)
    y = py0 + row * (ph_card + pgap)
    rect(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, pw, ph_card, fill=WHITE, line=LINE, line_w=0.75, radius=0.12)
    slide.shapes.add_picture(os.path.join(ICON_DIR, ico + ".png"),
                             Cm(x + 0.3), Cm(y + (ph_card - 0.95) / 2), Cm(0.95), Cm(0.95))
    txt(x + 1.25, y + 0.3, pw - 1.55, ph_card - 0.5, [[(text, 11, INK, True)]],
        line_spacing=1.35, anchor=MSO_ANCHOR.MIDDLE)

# ============ ④ メリット（解決） ============
mer_title_y = py0 + 2 * (ph_card + pgap) + 0.35
txt(1.2, mer_title_y, 17, 0.6,
    [[("そのお悩み、Takatsuki BASE が解決します", 12, GREEN, True, 1.0), ("   5 POINTS", 9, ORANGE, True, 1.5)]])
rect(MSO_SHAPE.RECTANGLE, 1.2, mer_title_y + 0.65, 18.6, 0.02, fill=LINE)

merits = [
    ("m_pin",    "地域の人に知ってもらえる", "高槻に住む人・働く人へお店を届ける"),
    ("m_send",   "SNSの外の層へ届く", "SNSを見ない世代・新規のお客様にも"),
    ("m_search", "Google検索から来店", "「高槻 ◯◯」で見つけてもらえる"),
    ("m_tag",    "イベント・セール告知", "新メニューや催しを自由に発信"),
    ("m_brief",  "求人も掲載できる", "スタッフ募集・採用にも役立つ"),
]
mcols = 5
mgap = 0.35
mw = (18.6 - mgap * (mcols - 1)) / mcols
mh_card = 4.1
mx0 = 1.2
mlist_y = mer_title_y + 1.0
for i, (ico, title, desc) in enumerate(merits):
    x = mx0 + i * (mw + mgap)
    rect(MSO_SHAPE.ROUNDED_RECTANGLE, x, mlist_y, mw, mh_card, fill=WHITE, line=LINE, line_w=0.75, radius=0.12)
    rect(MSO_SHAPE.RECTANGLE, x, mlist_y, mw, 0.1, fill=ORANGE)  # 上アクセント
    rect(MSO_SHAPE.ROUNDED_RECTANGLE, x + (mw - 1.2) / 2, mlist_y + 0.45, 1.2, 1.2,
         fill=RGBColor(0xe6, 0xf1, 0xeb), radius=0.28)
    slide.shapes.add_picture(os.path.join(ICON_DIR, ico + ".png"),
                             Cm(x + mw / 2 - 0.35), Cm(mlist_y + 0.7), Cm(0.7), Cm(0.7))
    txt(x + 0.15, mlist_y + 1.85, mw - 0.3, 1.0, [[(title, 11, INK, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.15)
    txt(x + 0.15, mlist_y + 2.85, mw - 0.3, 1.1, [[(desc, 8.5, SUB, False)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.25)

# ============ ⑤ CTA ＋ 連絡先 ＋ QR ============
cta_y = 20.0
rect(MSO_SHAPE.RECTANGLE, 0, cta_y, 21, 8.55, fill=GREEN)
txt(1.2, cta_y + 0.6, 14, 2.0, [
    [("まずは", 19, WHITE, True), ("無料", 19, PEACH, True), ("で、", 19, WHITE, True)],
    [("掲載の相談から始めませんか？", 19, WHITE, True)],
], line_spacing=1.3)
txt(1.2, cta_y + 3.0, 13.2, 1.6,
    [[("資料作成や掲載内容の整理もこちらでサポートします。「話だけ聞いてみたい」も大歓迎です。お気軽にご相談ください。",
       11, RGBColor(0xdf, 0xe9, 0xe2), False)]], line_spacing=1.45)
# 連絡先（2x2）
contacts = [
    ("運営", "株式会社RaidQ"), ("WEB", "takatsuki-base-portal.vercel.app"),
    ("__mail__", "info@raidq.jp"),
]
cy = cta_y + 5.0
for i, (k, v) in enumerate(contacts):
    col = i % 2; row = i // 2
    x = 1.2 + col * 7.0
    y = cy + row * 0.95
    if k == "__mail__":
        slide.shapes.add_picture(os.path.join(ICON_DIR, "mail.png"),
                                 Cm(x + 0.15), Cm(y + 0.05), Cm(0.55), Cm(0.55))
    else:
        txt(x, y, 0.9, 0.6, [[(k, 10.5, ORANGE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(x + 1.0, y, 6.0, 0.6, [[(v, 11, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
# QR
qr_box_x, qr_box_y, qr_box_w = 15.0, cta_y + 1.0, 4.6
rect(MSO_SHAPE.ROUNDED_RECTANGLE, qr_box_x, qr_box_y, qr_box_w, 5.4, fill=WHITE, radius=0.1)
slide.shapes.add_picture("LINEのQRコード.png", Cm(qr_box_x + 0.55), Cm(qr_box_y + 0.5), Cm(3.5), Cm(3.5))
txt(qr_box_x, qr_box_y + 4.25, qr_box_w, 0.7, [[("LINEで気軽にご相談", 9.5, GREEN, True)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============ フッター ============
rect(MSO_SHAPE.RECTANGLE, 0, 28.55, 21, 1.15, fill=GREEN_D)
txt(1.2, 28.55, 12, 1.15, [[("Takatsuki BASE ｜ 高槻のお店・人・イベントをつなぐ地域情報ポータル",
    8.5, RGBColor(0xc6, 0xd4, 0xcc), False)]], anchor=MSO_ANCHOR.MIDDLE)
txt(11.5, 28.55, 8.3, 1.15, [[("※ロゴは仮の表記です", 8.5, RGBColor(0xc6, 0xd4, 0xcc), False)]],
    align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

prs.save("Takatsuki_BASE_営業1枚資料.pptx")
print("saved: Takatsuki_BASE_営業1枚資料.pptx")
