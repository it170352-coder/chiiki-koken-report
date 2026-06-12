# -*- coding: utf-8 -*-
"""資料の画像.png（10枚一覧）を10スライドに分割して配置した PPTX を生成する。

元画像は低解像度（1536x1024）のため、SCALE 倍に高品質拡大（LANCZOS）し、
輪郭シャープ化（UnsharpMask）を施してから切り出す。これで投影・印刷時の
にじみを可能な範囲で軽減する（※元データに無い精細さは復元できない）。

16:9 スライドに、画像のアスペクト比を保ったまま中央配置する（白背景）。
"""
import os
from pptx import Presentation
from pptx.util import Inches as In
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image, ImageFilter

ROOT = "/Users/ishida/Desktop/claude code"
CELLS = os.path.join(ROOT, "slice_cells")
WHITE = RGBColor(0xff, 0xff, 0xff)
SCALE = 3  # 拡大倍率（1536 -> 4608px幅、各スライド約1500px幅）

SW, SH = 13.333, 7.5

# --- 元画像を高品質拡大 + シャープ化 ---
src = Image.open(os.path.join(ROOT, "資料の画像.png")).convert("RGB")
W, H = src.size
src_hi = src.resize((W * SCALE, H * SCALE), Image.LANCZOS)
src_hi = src_hi.filter(ImageFilter.UnsharpMask(radius=2.2, percent=130, threshold=2))

# --- グリッド座標（等倍）。切り出し時に SCALE 倍する ---
X = [3, 497, 1024, 1534]
ROWS = {1: (6, 272), 2: (304, 534), 3: (560, 767)}


def crop_hi(box):
    l, t, r, b = box
    return src_hi.crop((l * SCALE, t * SCALE, r * SCALE, b * SCALE))


os.makedirs(CELLS, exist_ok=True)
cell_paths = {}
n = 1
for ri in (1, 2, 3):
    y0, y1 = ROWS[ri]
    for ci in range(3):
        img = crop_hi((X[ci], y0, X[ci + 1], y1))
        p = os.path.join(CELLS, "cell%02d.png" % n)
        img.save(p)
        cell_paths[n] = p
        n += 1
# スライド10の2分割
closing = crop_hi((3, 793, 582, 1018))
cta = crop_hi((585, 793, 1534, 1018))
cp = os.path.join(CELLS, "cell10a_closing.png"); closing.save(cp)
tp = os.path.join(CELLS, "cell10b_cta.png"); cta.save(tp)

# --- PPTX 組み立て ---
prs = Presentation()
prs.slide_width = In(SW)
prs.slide_height = In(SH)
BLANK = prs.slide_layouts[6]


def white_bg(s):
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, In(SW), In(SH))
    bg.shadow.inherit = False
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()


def place_fit(s, path, area_left, area_top, area_w, area_h, valign="center"):
    iw, ih = Image.open(path).size
    disp_w = area_w
    disp_h = area_w * ih / iw
    if disp_h > area_h:
        disp_h = area_h
        disp_w = area_h * iw / ih
    left = area_left + (area_w - disp_w) / 2
    if valign == "top":
        top = area_top
    elif valign == "bottom":
        top = area_top + (area_h - disp_h)
    else:
        top = area_top + (area_h - disp_h) / 2
    s.shapes.add_picture(path, In(left), In(top), In(disp_w), In(disp_h))


# スライド1〜9
for k in range(1, 10):
    s = prs.slides.add_slide(BLANK)
    white_bg(s)
    place_fit(s, cell_paths[k], 0, 0, SW, SH)

# スライド10（締め＋CTAを上下分割）
s = prs.slides.add_slide(BLANK)
white_bg(s)
margin = 0.55
gap = 0.45
half_h = (SH - 2 * margin - gap) * 0.5
place_fit(s, cp, 0, margin, SW, half_h, valign="bottom")
place_fit(s, tp, 0, margin + half_h + gap, SW, half_h, valign="top")

out = os.path.join(ROOT, "03_資料スライド", "Takatsuki_BASE_営業資料_画像版.pptx")
prs.save(out)
print("saved:", out, "slides=", len(prs.slides._sldIdLst), "scale=", SCALE)
